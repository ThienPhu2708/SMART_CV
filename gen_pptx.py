"""
Tạo SmartCV_Nhom14.pptx — không xóa slide, chỉ modify in-place + thêm mới.
Tránh duplicate ZIP entry gây lỗi PowerPoint.
"""
import sys, copy, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

TEMPLATE = r'D:\Nam3_DH\DeepLearning\Do_An_Deeplearning.pptx'
OUTPUT   = r'D:\Nam3_DH\DeepLearning\SmartCV_Nhom14.pptx'

BG     = RGBColor(0xFF, 0xFD, 0xF6)
TITLE  = RGBColor(0x29, 0x29, 0x29)
BODY   = RGBColor(0x29, 0x29, 0x29)
BLUE   = RGBColor(0x2F, 0x5F, 0x98)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xE8, 0xF0, 0xFE)
PBLUE  = RGBColor(0xF0, 0xF4, 0xFF)
NCOLOR = RGBColor(0x22, 0x22, 0x22)

TF = 'Public Sans Bold'
BF = 'Calibri'
NF = 'Tahoma'

# ─── helpers ──────────────────────────────────────────────────────────────────

def clear_shapes(slide, keep_keyword='Freeform'):
    """Xóa tất cả shapes trừ freeforms."""
    spTree = slide.shapes._spTree
    to_rm = [sh.element for sh in slide.shapes
              if keep_keyword not in sh.name]
    for el in to_rm:
        spTree.remove(el)

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, l, t, w, h, text, fn, sz, bold=False,
       color=TITLE, align=PP_ALIGN.LEFT, wrap=True):
    """Thêm text box đơn giản."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.name  = fn
    r.font.size  = Pt(sz)
    r.font.bold  = bold
    r.font.color.rgb = color

def bullets(slide, l, t, w, h, items, fn=BF, sz=24, color=BODY, char='•'):
    """Thêm text box nhiều dòng bullet."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5)
        r = p.add_run()
        r.text = (f'{char} {item}' if char and item.strip() else item)
        r.font.name  = fn
        r.font.size  = Pt(sz)
        r.font.color.rgb = color

def sep(slide, l=1.8, t=1.85, w=16.2):
    """Đường kẻ ngang màu xanh."""
    ln = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE
    ln.line.fill.background()

def pagenum(slide, n):
    tb(slide, 18.88, 10.07, 0.6, 0.3, str(n), NF, 20, color=NCOLOR, align=PP_ALIGN.RIGHT)

def clone_freeforms(slide, src_slide):
    """Copy freeform decorative shapes từ src_slide, bao gồm image relationships."""
    rid_map = {}
    for rel in src_slide.part.rels.values():
        if 'image' in rel.reltype:
            new_rid = slide.part.relate_to(rel._target, rel.reltype)
            rid_map[rel.rId] = new_rid

    spTree = slide.shapes._spTree
    for elem in list(src_slide.shapes._spTree)[2:]:
        name = ''
        cNvPr = elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
        if cNvPr is None:
            # try drawingml
            cNvPr = elem.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr')
        if cNvPr is not None:
            name = cNvPr.get('name', '')
        if 'Freeform' not in name:
            continue
        new_el = copy.deepcopy(elem)
        xml_b = etree.tostring(new_el)
        for old, new in rid_map.items():
            xml_b = xml_b.replace(f'r:embed="{old}"'.encode(),
                                   f'r:embed="{new}"'.encode())
        spTree.append(etree.fromstring(xml_b))

# ─── slide builders ────────────────────────────────────────────────────────────

def chapter_header(slide, src_ch, title, n):
    set_bg(slide)
    clone_freeforms(slide, src_ch)
    tb(slide, 2.5, 4.0, 15, 3.2, title, TF, 52, bold=True,
       color=TITLE, align=PP_ALIGN.CENTER)
    pagenum(slide, n)

def content(slide, src_ct, title, items, n, sz=24):
    set_bg(slide)
    clone_freeforms(slide, src_ct)
    tb(slide, 1.8, 0.6, 16.2, 1.1, title, TF, 40, bold=True, color=TITLE)
    sep(slide)
    bullets(slide, 1.8, 2.0, 16.2, 8.6, items, sz=sz)
    pagenum(slide, n)

def two_col(slide, src_ct, title, lt, li, rt, ri, n):
    set_bg(slide)
    clone_freeforms(slide, src_ct)
    tb(slide, 1.8, 0.6, 16.2, 1.1, title, TF, 40, bold=True, color=TITLE)
    sep(slide)
    tb(slide, 1.8, 2.0, 7.8, 0.7, lt, TF, 26, bold=True, color=BLUE)
    bullets(slide, 1.8, 2.8, 7.8, 7.8, li, sz=22)
    div = slide.shapes.add_shape(1, Inches(9.9), Inches(2.0), Inches(0.03), Inches(8.5))
    div.fill.solid(); div.fill.fore_color.rgb = BLUE
    div.line.fill.background()
    tb(slide, 10.1, 2.0, 7.9, 0.7, rt, TF, 26, bold=True, color=BLUE)
    bullets(slide, 10.1, 2.8, 7.9, 7.8, ri, sz=22)
    pagenum(slide, n)

def table_slide(slide, src_ct, title, headers, rows, n):
    set_bg(slide)
    clone_freeforms(slide, src_ct)
    tb(slide, 1.8, 0.6, 16.2, 1.1, title, TF, 40, bold=True, color=TITLE)
    sep(slide)
    nc, nr = len(headers), len(rows)+1
    row_h = min(0.85, 7.8 / nr)
    tbl = slide.shapes.add_table(nr, nc,
          Inches(1.8), Inches(2.0),
          Inches(16.2), Inches(row_h * nr)).table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = BLUE
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = BF; r.font.size = Pt(18)
        r.font.bold = True; r.font.color.rgb = WHITE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i+1, j)
            c.text = str(val)
            if i % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = LGRAY
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.runs[0]
            r.font.name = BF; r.font.size = Pt(17)
            r.font.color.rgb = BODY
    pagenum(slide, n)

def img_slide(slide, src_ct, title, caption, n):
    set_bg(slide)
    clone_freeforms(slide, src_ct)
    tb(slide, 1.8, 0.6, 16.2, 1.1, title, TF, 40, bold=True, color=TITLE)
    sep(slide)
    box = slide.shapes.add_shape(1, Inches(2.8), Inches(2.1), Inches(14.4), Inches(7.7))
    box.fill.solid(); box.fill.fore_color.rgb = PBLUE
    box.line.color.rgb = BLUE; box.line.width = Pt(1.5)
    tb(slide, 2.8, 5.0, 14.4, 2.0, caption, BF, 22,
       color=BLUE, align=PP_ALIGN.CENTER)
    pagenum(slide, n)

# ─── main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation(TEMPLATE)
    tpl_ch  = prs.slides[2]   # slide 3 — chapter header freeforms
    tpl_ct  = prs.slides[3]   # slide 4 — content freeforms

    # ── Slide 1: Title — sửa "Nhóm 7" → "Nhóm 14" ──────────────────────
    for sh in prs.slides[0].shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if 'Nhóm 7' in run.text:
                        run.text = 'Nhóm 14'

    # ── Slide 2: Members — đã đúng ───────────────────────────────────────
    # ── Slide 3: Ch1 header — đã đúng ────────────────────────────────────

    # Helper: lấy slide hiện có (index 3-16) hoặc thêm mới
    blank = prs.slide_layouts[6]
    def get_slide(idx):
        """Lấy slide idx nếu có (và clear), hoặc thêm mới."""
        if idx < len(prs.slides):
            sl = prs.slides[idx]
            clear_shapes(sl)
            set_bg(sl)
            return sl
        return prs.slides.add_slide(blank)

    # ── Slide 4 (idx 3): Bối cảnh & Vấn đề ─────────────────────────────
    s = get_slide(3)
    content(s, tpl_ct, 'BỐI CẢNH VÀ VẤN ĐỀ', [
        'Thị trường lao động 4.0 bùng nổ → số lượng hồ sơ tăng đột biến',
        'HR phải xử lý hàng trăm CV thủ công — tốn thời gian, dễ sai sót',
        'Đánh giá cảm tính, thiếu khách quan → bỏ lọt ứng viên tốt',
        'Keyword matching đơn giản không hiểu ngữ cảnh (TF-IDF)',
        '→ Cần hệ thống tự động hóa giai đoạn sơ loại hồ sơ',
    ], 4)

    # ── Slide 5 (idx 4): Bài toán & Giải pháp ───────────────────────────
    s = get_slide(4)
    two_col(s, tpl_ct, 'BÀI TOÁN & GIẢI PHÁP',
        'BÀI TOÁN',
        ['Sàng lọc CV tự động thay thế con người giai đoạn sơ loại',
         'Phân loại ngành nghề phù hợp với vị trí tuyển dụng',
         'Đánh giá điều kiện bắt buộc theo quy tắc nghiệp vụ HR'],
        'GIẢI PHÁP ĐỀ XUẤT',
        ['Pipeline 4 tầng: OCR/PDF → BERT → MLP → Logic Gates',
         'Hỗ trợ PDF, ảnh scan (EasyOCR)',
         'Hỗ trợ Tiếng Anh & Tiếng Việt (multilingual)',
         'Logic Gates: AND / OR / NOT / XOR do HR cấu hình'],
        5)

    # ── Slide 6 (idx 5): Ch2 header ──────────────────────────────────────
    s = get_slide(5)
    chapter_header(s, tpl_ch, 'CHƯƠNG 2: PHÂN TÍCH ĐỀ TÀI', 6)

    # ── Slide 7 (idx 6): Phân tích yêu cầu ──────────────────────────────
    s = get_slide(6)
    content(s, tpl_ct, 'PHÂN TÍCH YÊU CẦU', [
        'Mục tiêu: Hệ thống thông minh thay thế con người sơ loại hồ sơ',
        'Chức năng chính: Trích xuất text, đánh giá logic, dự báo độ phù hợp',
        'Vai trò: Xác định đặc trưng cần học, định hướng toàn bộ pipeline',
        'Thách thức: CV đa dạng bố cục, TF-IDF không hiểu ngữ cảnh, cân bằng Logic vs ML',
        'Định hướng: Pipeline 4 tầng BERT + MLP + Logic Gates',
    ], 7)

    # ── Slide 8 (idx 7): Yêu cầu chức năng ──────────────────────────────
    s = get_slide(7)
    content(s, tpl_ct, 'YÊU CẦU CHỨC NĂNG', [
        'HR Setup — Cấu hình vị trí tuyển dụng, kỹ năng yêu cầu, điều kiện Logic Gates',
        'Input Management — Tiếp nhận PDF / JPG / PNG, kiểm tra định dạng file',
        'Core Processing — OCR → Phát hiện ngôn ngữ → Dịch → BERT → MLP → Logic',
        'Result Visualization — Hiển thị ĐẠT/LOẠI, confidence%, sắp xếp, xuất CSV',
        'History Tracking — Lưu trữ và tra cứu lịch sử các đợt sàng lọc',
    ], 8)

    # ── Slide 9 (idx 8): Ch3 header ──────────────────────────────────────
    s = get_slide(8)
    chapter_header(s, tpl_ch, 'CHƯƠNG 3: THIẾT KẾ ỨNG DỤNG', 9)

    # ── Slide 10 (idx 9): ANNs/MLP ───────────────────────────────────────
    s = get_slide(9)
    content(s, tpl_ct, 'MẠNG NƠ-RON NHÂN TẠO (ANNs / MLP)', [
        'ANN: mô hình học dựa trên hệ thần kinh sinh học — học mối quan hệ phi tuyến',
        'MLP (Multi-Layer Perceptron): mạng truyền thẳng Feed Forward',
        'Kiến trúc cơ bản: Input Layer → Hidden Layers → Output Layer',
        'Công thức neuron: y = f(Σ wᵢxᵢ + b)',
        'Hàm kích hoạt: ReLU  f(x) = max(0, x)',
        'Output layer: Softmax → xác suất thuộc từng ngành',
    ], 10)

    # ── Slide 11 (idx 10): BERT Encoder ──────────────────────────────────
    s = get_slide(10)
    content(s, tpl_ct, 'TẦNG 2: BERT ENCODER (FROZEN)', [
        'Model: paraphrase-multilingual-MiniLM-L12-v2 (Sentence-Transformer)',
        'Chuyển văn bản CV → vector ngữ nghĩa 384 chiều',
        'Hiểu ngữ cảnh: "SEO" trong CV IT  ≠  "SEO" trong CV Digital Media',
        'Khác TF-IDF (chỉ đếm từ) → BERT hiểu toàn bộ câu văn',
        'Trọng số FROZEN hoàn toàn — chỉ train MLP Head phía sau',
        'Hỗ trợ Tiếng Anh + Tiếng Việt (multilingual model)',
    ], 11)

    # ── Slide 12 (idx 11): MLP Head ──────────────────────────────────────
    s = get_slide(11)
    content(s, tpl_ct, 'TẦNG 3: MLP CLASSIFICATION HEAD', [
        'Input Layer:    384 neuron  (BERT embedding vector)',
        'Hidden Layer 1: 256 neuron  →  BatchNorm1d  →  ReLU  →  Dropout(0.45)',
        'Hidden Layer 2: 128 neuron  →  BatchNorm1d  →  ReLU  →  Dropout(0.35)',
        'Output Layer:     6 ngành   →  Softmax  →  confidence%',
        '',
        'Tổng tham số: ~132.000  (gọn nhẹ, tránh overfitting ~1.197 mẫu)',
        'BatchNorm: ổn định training  |  Dropout: ngăn ghi nhớ dữ liệu',
    ], 12, sz=24)

    # ── Slide 13 (idx 12): Logic Gates AND/OR ────────────────────────────
    s = get_slide(12)
    two_col(s, tpl_ct, 'TẦNG 4: LOGIC GATES',
        'AND Gate — Bắt buộc TẤT CẢ',
        ['Ví dụ: Python AND SQL AND Docker',
         'Ứng viên chỉ ĐẠT nếu có đủ tất cả kỹ năng',
         '0 AND 0 = 0   |   0 AND 1 = 0',
         '1 AND 0 = 0   |   1 AND 1 = 1'],
        'OR Gate — Ít nhất 1 điều kiện',
        ['Ví dụ: ReactJS OR VueJS OR Angular',
         'Ứng viên ĐẠT nếu có ít nhất 1 kỹ năng',
         '0 OR 0 = 0   |   0 OR 1 = 1',
         '1 OR 0 = 1   |   1 OR 1 = 1'],
        13)

    # ── Slide 14 (idx 13): Logic Gates NOT/XOR ───────────────────────────
    s = get_slide(13)
    two_col(s, tpl_ct, 'TẦNG 4: LOGIC GATES (tiếp theo)',
        'NOT Gate — Loại trừ',
        ['Ví dụ: NOT "fresher"  |  NOT "intern"',
         'LOẠI nếu CV chứa từ khóa cấm',
         'NOT 0 = 1   |   NOT 1 = 0',
         'Mô phỏng tiêu chí bắt buộc của HR'],
        'XOR Gate — Chọn đúng 1',
        ['Ví dụ: Frontend XOR Backend',
         'ĐẠT nếu đúng 1 điều kiện, không cả hai',
         '0 XOR 0 = 0   |   0 XOR 1 = 1',
         '1 XOR 0 = 1   |   1 XOR 1 = 0'],
        14)

    # ── Slide 15 (idx 14): Quy trình tổng thể ────────────────────────────
    s = get_slide(14)
    content(s, tpl_ct, 'QUY TRÌNH XỬ LÝ TỔNG THỂ', [
        '① Upload CV (PDF / JPG / PNG)',
        '② Trích xuất text: pdfplumber → PyPDF2 → EasyOCR (fallback)',
        '③ Phát hiện ngôn ngữ: ký tự tiếng Việt > 8% → dịch sang tiếng Anh',
        '④ Tiền xử lý: lowercase, xóa URL/email, chuẩn hóa viết tắt, Skills Boost, Lemmatization',
        '⑤ BERT Encoder: văn bản → vector 384 chiều',
        '⑥ MLP Head: 384 → 256 → 128 → 6 ngành → Softmax → confidence%',
        '⑦ Logic Gates: AND/OR/NOT/XOR → ĐẠT / LOẠI  +  cảnh báo nếu confidence ≥ 30%',
    ], 15, sz=23)

    # ── Slide 16 (idx 15): Phương pháp trích xuất text ───────────────────
    s = get_slide(15)
    table_slide(s, tpl_ct, 'PHƯƠNG PHÁP TRÍCH XUẤT TEXT',
        ['Tiêu chí', 'pdfplumber (ưu tiên)', 'PyPDF2 (dự phòng)', 'EasyOCR (fallback)'],
        [
            ['CV đa cột / sidebar',  'Tốt (sắp xếp tọa độ)', 'Dễ lẫn lộn cột',  'Không áp dụng'],
            ['CV đơn giản 1 cột',    'Tốt',                   'Tốt',              'Không áp dụng'],
            ['PDF scan / ảnh',       'Không hỗ trợ',          'Không hỗ trợ',     'Hỗ trợ (OCR)'],
            ['Tốc độ xử lý',         'Trung bình',             'Nhanh',            'Chậm'],
            ['Ngưỡng kích hoạt',     'Mặc định ưu tiên',      '< 50 ký tự',       '< 50 ký tự'],
        ],
        16)

    # ── Slide 17 (idx 16): Tiền xử lý ────────────────────────────────────
    s = get_slide(16)
    content(s, tpl_ct, 'PHƯƠNG PHÁP TIỀN XỬ LÝ DỮ LIỆU', [
        'Lowercase: Python → python, PYTHON → python',
        'Xóa URL, địa chỉ email, số điện thoại, ký tự đặc biệt',
        'Chuẩn hóa viết tắt: ML→machine learning, AI→artificial intelligence',
        'Skills Boosting: nhân đôi phần "Skills" → từ kỹ năng có trọng số cao hơn',
        'Lemmatization: đưa về dạng gốc (running→run, experiences→experience)',
        'Loại bỏ stop words tiếng Anh (NLTK)',
        'Phát hiện tiếng Việt (ngưỡng 8%) → dịch sang tiếng Anh tự động',
    ], 17, sz=23)

    # ═══ Từ đây thêm slides mới (idx 17+) — không gây duplicate ══════════

    def ns():
        return prs.slides.add_slide(blank)

    # ── Slide 18: Dataset nguồn gốc ──────────────────────────────────────
    s = ns()
    content(s, tpl_ct, 'NGUỒN GỐC VÀ THÔNG SỐ DATASET', [
        'Nguồn 1 — Kaggle UpdatedResumeDataSet: 2.484 hồ sơ / 24 ngành → lọc lấy 678 mẫu',
        'Nguồn 2 — Neuralframe Dataset: 9.544 mẫu cấu trúc → giới hạn tối đa 80 mẫu/ngành',
        'Nguồn 3 — Synthetic Data (tự sinh): 235 mẫu template-based generation, prefix SYN_',
        '',
        '6 NGÀNH PHÂN LOẠI:',
        'INFORMATION-TECHNOLOGY  |  BUSINESS-DEVELOPMENT  |  SALES',
        'DIGITAL-MEDIA  |  PUBLIC-RELATIONS  |  CONSULTANT',
    ], 18, sz=24)

    # ── Slide 19: Phân bổ dữ liệu ────────────────────────────────────────
    s = ns()
    table_slide(s, tpl_ct, 'CẤU TRÚC VÀ PHÂN BỔ DỮ LIỆU',
        ['Ngành', 'Kaggle', 'Neuralframe', 'Synthetic', 'Tổng'],
        [
            ['INFORMATION-TECHNOLOGY', '120', '80',  '0',   '200'],
            ['BUSINESS-DEVELOPMENT',   '120', '40',  '41',  '201'],
            ['SALES',                  '116', '56',  '28',  '200'],
            ['PUBLIC-RELATIONS',       '111', '80',  '5',   '196'],
            ['CONSULTANT',             '115', '28',  '57',  '200'],
            ['DIGITAL-MEDIA',          '96',  '0',   '104', '200'],
            ['TỔNG',                   '678', '284', '235', '~1.197'],
        ], 19)

    # ── Slide 20: Train/Val/Test ──────────────────────────────────────────
    s = ns()
    table_slide(s, tpl_ct, 'PHÂN CHIA DỮ LIỆU — TRAIN / VAL / TEST',
        ['Tập', 'Tỉ lệ', 'Mục đích'],
        [
            ['Train',      '70%', 'Huấn luyện trọng số model'],
            ['Validation', '15%', 'Early Stopping — chọn checkpoint tốt nhất'],
            ['Test',       '15%', 'Đánh giá cuối — chỉ dùng 1 lần (unbiased)'],
        ], 20)

    # ── Slide 21: Ch4 header ─────────────────────────────────────────────
    s = ns()
    chapter_header(s, tpl_ch, 'CHƯƠNG 4: THỰC HIỆN CÀI ĐẶT', 21)

    # ── Slide 22: Môi trường ─────────────────────────────────────────────
    s = ns()
    two_col(s, tpl_ct, 'CÔNG CỤ VÀ MÔI TRƯỜNG PHÁT TRIỂN',
        'Deep Learning & NLP',
        ['PyTorch ≥ 2.0',
         'Sentence-Transformers ≥ 2.2 (BERT)',
         'Scikit-learn ≥ 1.2',
         'NLTK (Lemmatization, Stopwords)'],
        'PDF / OCR / Web',
        ['pdfplumber ≥ 0.10   |   PyPDF2',
         'EasyOCR ≥ 1.7   |   opencv-python',
         'deep-translator   |   langdetect',
         'FastAPI ≥ 0.110   |   Uvicorn   |   Jinja2'],
        22)

    # ── Slide 23: MLP Baseline ────────────────────────────────────────────
    s = ns()
    content(s, tpl_ct, 'GIAI ĐOẠN 1: MLP BASELINE (TF-IDF)', [
        'Kiến trúc: TF-IDF (600 dim) → Linear(64) → BN → ReLU → Dropout(0.4)',
        '                            → Linear(32) → BN → ReLU → Dropout(0.3) → Linear(6)',
        'Tổng tham số: ~41.000',
        '',
        'Siêu tham số: Epochs=150 | Batch=64 | LR=0.001 | CosineAnnealingLR | EarlyStopping=8',
        '',
        'KẾT QUẢ:  Best Val F1 = 87.73%  |  Val Acc = 87.71%  |  Epoch thực tế = 17',
        'Test Accuracy = 85.56%  |  Test Macro F1 = 85.67%  |  Overfit gap = ~13.4% ⚠',
    ], 23, sz=23)

    # ── Slide 24: Nguyên nhân Overfit ─────────────────────────────────────
    s = ns()
    content(s, tpl_ct, 'NGUYÊN NHÂN OVERFIT — TF-IDF BAG-OF-WORDS', [
        'TF-IDF chỉ đếm tần suất từ — KHÔNG phân biệt ngữ cảnh',
        '',
        'Ví dụ điển hình:',
        '   "SEO" trong CV Frontend Developer (TypeScript / React)',
        '   ≡ "SEO" trong CV Digital Marketer  →  MLP nhầm ngành!',
        '',
        'Từ "management" xuất hiện ở CONSULTANT, BUSINESS-DEV, SALES → ranh giới mờ',
        'MLP buộc phải GHI NHỚ pattern bề mặt thay vì học đặc trưng thực sự',
        '',
        '→ GIẢI PHÁP: Thay TF-IDF bằng BERT Contextual Embedding',
    ], 24, sz=23)

    # ── Slide 25: BERT + MLP Pipeline ─────────────────────────────────────
    s = ns()
    content(s, tpl_ct, 'GIAI ĐOẠN 2: BERT + MLP (PIPELINE CHÍNH THỨC)', [
        'BERT Encoder frozen → 384-dim embedding (hiểu ngữ cảnh toàn câu)',
        'Classification head vẫn là MLP thuần túy (384 → 256 → 128 → 6)',
        '→ Đúng yêu cầu đề tài "Áp dụng MLP", cải thiện chất lượng đặc trưng đầu vào',
        '',
        'Siêu tham số:  Epochs=80  |  Batch=32  |  LR=2e-4  |  Optimizer=AdamW',
        'Weight_decay=8e-3  |  ReduceLROnPlateau (factor=0.5, patience=5)',
        'Early stopping: patience=15  |  LabelSmoothing=0.1  |  ClassWeights',
        '',
        'Tại sao ReduceLROnPlateau thay CosineAnnealingLR?',
        '→ CosineAnnealing gây dao động Val F1 ±3% → ReduceLROnPlateau hội tụ mượt hơn',
    ], 25, sz=22)

    # ── Slide 26: Training MLP chart ──────────────────────────────────────
    s = ns()
    img_slide(s, tpl_ct, 'ĐÁNH GIÁ HUẤN LUYỆN — MLP BASELINE',
        '[ Chèn hình: Models/training_curves.png ]\n\nLoss / Accuracy / Val F1 theo Epoch\nEpoch 17  |  Train Acc ~99%  vs  Val Acc ~87%  →  Overfit rõ rệt', 26)

    # ── Slide 27: Training BERT chart ─────────────────────────────────────
    s = ns()
    img_slide(s, tpl_ct, 'ĐÁNH GIÁ HUẤN LUYỆN — BERT + MLP',
        '[ Chèn hình: Models/bert_training_curves.png ]\n\nLoss / Accuracy / Val F1 theo Epoch\nTrain ≈ Val xuyên suốt  →  Overfit gần như biến mất khi dùng BERT', 27)

    # ── Slide 28: So sánh hai mô hình ─────────────────────────────────────
    s = ns()
    table_slide(s, tpl_ct, 'ĐÁNH GIÁ VÀ SO SÁNH HAI MÔ HÌNH',
        ['Chỉ số', 'MLP Baseline (TF-IDF)', 'BERT + MLP ✓'],
        [
            ['Val Accuracy',  '87.71%',      '91.62%'],
            ['Val Macro F1',  '87.73%',      '91.54%'],
            ['Test Accuracy', '85.56%',      '88.89%'],
            ['Test Macro F1', '85.67%',      '88.84%'],
            ['Overfit gap',   '~13.4% ⚠',   '~2.1% ✓'],
            ['Embedding',     'Bag-of-words (600 dim)', 'Contextual (384 dim)'],
        ], 28)

    # ── Slide 29: Confusion Matrix ─────────────────────────────────────────
    s = ns()
    img_slide(s, tpl_ct, 'CONFUSION MATRIX — BERT + MLP',
        '[ Chèn hình: Models/confusion_matrix.png ]\n\nMa trận nhầm lẫn trên tập Test 15%', 29)

    # ── Slide 30: Phân tích confusion matrix ──────────────────────────────
    s = ns()
    content(s, tpl_ct, 'PHÂN TÍCH KẾT QUẢ', [
        'Đường chéo đậm → mô hình phân loại chính xác đa số ngành',
        'INFORMATION-TECHNOLOGY: ít nhầm nhất — từ vựng kỹ thuật đặc trưng',
        'CONSULTANT ↔ BUSINESS-DEVELOPMENT: hay nhầm (CV tư vấn đa lĩnh vực)',
        'PUBLIC-RELATIONS ↔ DIGITAL-MEDIA: nhầm do cùng liên quan truyền thông',
        'SALES ↔ BUSINESS-DEVELOPMENT: từ vựng doanh thu, khách hàng tương đồng',
        '',
        '✓ KHÔNG có nhầm lẫn chéo xa (IT không nhầm sang SALES)',
    ], 30, sz=24)

    # ── Slide 31: Out-of-sample ────────────────────────────────────────────
    s = ns()
    table_slide(s, tpl_ct, 'ĐÁNH GIÁ TRÊN DỮ LIỆU THỰC TẾ (60 CV)',
        ['Ngành', 'Kết quả', 'Accuracy'],
        [
            ['BUSINESS-DEVELOPMENT',   '10/10', '100%'],
            ['INFORMATION-TECHNOLOGY', '9/10',  '90%'],
            ['DIGITAL-MEDIA',          '9/10',  '90%'],
            ['PUBLIC-RELATIONS',       '7/10',  '70%'],
            ['SALES',                  '6/10',  '60%'],
            ['CONSULTANT',             '5/10',  '50%'],
            ['TỔNG',                   '46/60', '76.7%'],
        ], 31)

    # ── Slide 32: Confidence ───────────────────────────────────────────────
    s = ns()
    two_col(s, tpl_ct, 'PHÂN TÍCH ĐỘ TIN CẬY (CONFIDENCE)',
        'Dự đoán ĐÚNG (46 CV)',
        ['Confidence tập trung 70 – 100%',
         'Đỉnh tại 80 – 90%',
         'Mô hình đưa ra quyết định có cơ sở',
         'Không đoán mò — tin cậy cao'],
        'Dự đoán SAI (14 CV)',
        ['Confidence tập trung 30 – 60%',
         'Đỉnh tại 35 – 45%',
         'Mô hình "nhận thức được" sự không chắc',
         'Logic Gates: cảnh báo HR nếu confidence ≥ 30%'],
        32)

    # ── Slides 33–37: Giao diện ───────────────────────────────────────────
    for title, fig, n in [
        ('GIAO DIỆN DASHBOARD',          'Hình 4.4 — Dashboard tổng quan', 33),
        ('GIAO DIỆN THIẾT LẬP JOB',      'Hình 4.5 — Tạo Job + cấu hình Logic Gates', 34),
        ('GIAO DIỆN UPLOAD CV',          'Hình 4.6 — Upload PDF/JPG/PNG', 35),
        ('GIAO DIỆN KẾT QUẢ SÀNG LỌC',  'Hình 4.7 — Danh sách ứng viên + ĐẠT/LOẠI', 36),
        ('CHI TIẾT CV ĐƯỢC ĐÁNH GIÁ',   'Hình 4.10 — Text CV + ngành + kết quả Logic Gate', 37),
    ]:
        img_slide(ns(), tpl_ct, title, f'[ Chèn screenshot: {fig} ]', n)

    # ── Slide 38: Ch5 header ───────────────────────────────────────────────
    s = ns()
    chapter_header(s, tpl_ch, 'CHƯƠNG 5: KẾT LUẬN VÀ HƯỚNG PHÁT TRIỂN', 38)

    # ── Slide 39: Kết luận ────────────────────────────────────────────────
    s = ns()
    content(s, tpl_ct, 'KẾT LUẬN', [
        'Xây dựng thành công hệ thống sàng lọc CV tự động 4 tầng',
        'BERT + MLP đạt Test Accuracy 88.89%  |  Overfit gap chỉ 2.1%',
        'MLP Baseline (TF-IDF) giữ lại để minh chứng hiệu quả cải tiến',
        'Web app FastAPI: PDF, ảnh scan, Tiếng Việt, lịch sử, xuất CSV',
        '',
        'KẾT LUẬN QUAN TRỌNG:',
        'Nguyên nhân overfit là CHẤT LƯỢNG ĐẶC TRƯNG (TF-IDF), không phải kiến trúc MLP',
        '→ Cùng MLP: thay TF-IDF bằng BERT  →  Test Acc 85.56% → 88.89%',
    ], 39, sz=24)

    # ── Slide 40: Hướng phát triển ────────────────────────────────────────
    s = ns()
    content(s, tpl_ct, 'HƯỚNG PHÁT TRIỂN', [
        'Mở rộng dataset: thêm ngành mới (Healthcare, Education, Finance)',
        'Cải thiện mô hình: Fine-tuning toàn bộ BERT encoder (+2–5% F1 kỳ vọng)',
        'Explainable AI: SHAP / LIME → HR hiểu lý do quyết định Đạt/Loại',
        'Triển khai cloud: Docker + AWS/GCP phục vụ nhiều người dùng',
        'Phân tích đa tiêu chí: số năm kinh nghiệm, bằng cấp, chứng chỉ',
        'Tích hợp API với các hệ thống ATS (Applicant Tracking System)',
    ], 40, sz=24)

    # ── Slide 41: Thank You ───────────────────────────────────────────────
    s = ns()
    chapter_header(s, tpl_ch,
        'THANK YOU!\n\nNhóm 14 — HUIT Deep Learning\nGVHD: TS. Nguyễn Thanh Long',
        41)

    # ── Lưu ──────────────────────────────────────────────────────────────
    prs.save(OUTPUT)
    print(f'Done — {len(prs.slides)} slides → {OUTPUT}')

if __name__ == '__main__':
    main()
